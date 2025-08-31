from pptx import Presentation
from typing import List, Optional
import logging
from .template_manager import template_manager, AspectRatio, TemplateCategory


class PPTGenerator:
    def __init__(self, output_filename: str, font_name: Optional[str] = None, 
                 template_id: str = "default_16_9", aspect_ratio: AspectRatio = None):
        """Initialize with template, output filename, and optional font."""
        self.output_filename = output_filename
        self.template_id = template_id
        
        # Load the template presentation using template manager
        self.template = template_manager.load_template(template_id)
        if not self.template:
            # Fallback to loading from file system
            try:
                self.template = Presentation("template.pptx")
                logging.warning(f"Template '{template_id}' not found, using fallback template.pptx")
            except:
                self.template = Presentation("generate_content/template.pptx")
                logging.warning(f"Using generate_content/template.pptx as fallback")
        
        # Get template metadata
        self.template_metadata = template_manager.get_template_metadata(template_id)
        
        self.presentation = Presentation()
        
        # Set font with error handling
        self.font_name = self._validate_font(font_name)
        
        # Set dimensions based on aspect ratio or template
        if aspect_ratio:
            width, height = template_manager.get_aspect_ratio_dimensions(aspect_ratio)
            self.presentation.slide_width = width
            self.presentation.slide_height = height
            logging.info(f"Set custom aspect ratio: {aspect_ratio.value}")
        else:
            # Copy slide dimensions from template to maintain aspect ratio
            self.presentation.slide_width = self.template.slide_width
            self.presentation.slide_height = self.template.slide_height
        
        # Clear any default slides from new presentation
        while len(self.presentation.slides) > 0:
            rId = self.presentation.slides._sldIdLst[0].rId
            self.presentation.part.drop_rel(rId)
            del self.presentation.slides._sldIdLst[0]
        
        logging.info(f"PPTGenerator initialized with template: {template_id}, aspect ratio: {aspect_ratio}")
    
    def _validate_font(self, font_name: Optional[str]) -> str:
        """Validate font availability and return appropriate font name."""
        # Default fonts that are commonly available across systems
        default_fonts = ["Arial", "Calibri", "Times New Roman", "Helvetica"]
        
        # If no font specified, use first default
        if not font_name:
            return default_fonts[0]
        
        # List of commonly available fonts
        available_fonts = [
            "Arial", "Arial Black", "Arial Narrow", "Arial Unicode MS",
            "Calibri", "Calibri Light", "Cambria", "Candara",
            "Century Gothic", "Comic Sans MS", "Consolas", "Constantia",
            "Corbel", "Courier New", "Franklin Gothic Medium",
            "Garamond", "Georgia", "Helvetica", "Impact",
            "Lucida Console", "Lucida Sans Unicode", "Microsoft Sans Serif",
            "Palatino Linotype", "Segoe UI", "Tahoma", "Times New Roman",
            "Trebuchet MS", "Verdana"
        ]
        
        try:
            # Check if the requested font is in our available list
            if font_name in available_fonts:
                logging.info(f"Using requested font: {font_name}")
                return font_name
            else:
                # Font not found, use default
                default_font = default_fonts[0]
                logging.warning(f"Font '{font_name}' not found. Using default font: {default_font}")
                return default_font
                
        except Exception as e:
            # Any error in font validation, use default
            default_font = default_fonts[0]
            logging.error(f"Error validating font '{font_name}': {e}. Using default font: {default_font}")
            return default_font
    
    def _apply_font_to_shape(self, shape):
        """Apply the selected font to a text shape."""
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font_name
        except Exception as e:
            logging.warning(f"Could not apply font '{self.font_name}' to shape: {e}")
            # Continue without applying font if there's an error
    
    def _add_citations_to_slide(self, slide, citations: List[str]):
        """Add citations as a text box at the bottom of the slide."""
        if not citations:
            return
        
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # Create a text box for citations at the bottom of the slide
            left = Inches(0.5)
            top = Inches(6.5)  # Near bottom of slide
            width = Inches(9)
            height = Inches(1)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Add "References:" header
            p = text_frame.paragraphs[0]
            p.text = "References:"
            p.font.size = Pt(10)
            p.font.bold = True
            if self.font_name:
                p.font.name = self.font_name
            
            # Add each citation
            for citation in citations:
                p = text_frame.add_paragraph()
                p.text = f"• {citation}"
                p.font.size = Pt(8)
                if self.font_name:
                    p.font.name = self.font_name
                p.level = 0
            
            # Set text alignment
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                
        except Exception as e:
            logging.warning(f"Could not add citations to slide: {e}")
            # Continue without adding citations if there's an error
    
    def add_title_slide(self, title_text: str, citations: List[str] = None):
        """Add title slide (slide 1 from template) with title text."""
        # Copy slide 0 (title slide) from template
        template_slide = self.template.slides[0]
        slide_layout = template_slide.slide_layout
        new_slide = self.presentation.slides.add_slide(slide_layout)
        
        # Copy all shapes from template slide
        for shape in template_slide.shapes:
            if not shape.is_placeholder:
                # Copy non-placeholder shapes (background elements)
                continue
        
        # Find and update the title textbox
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = title_text
                self._apply_font_to_shape(shape)
                break
        
        # Add citations if provided
        self._add_citations_to_slide(new_slide, citations)
    
    def add_bullet_slide(self, heading_text: str, bullet_points: List[str], citations: List[str] = None):
        """Add bullet slide (slide 2 from template) with heading and bullet points."""
        # Copy slide 1 (bullet slide) from template
        template_slide = self.template.slides[1]
        slide_layout = template_slide.slide_layout
        new_slide = self.presentation.slides.add_slide(slide_layout)
        
        # Find textboxes and update content
        textboxes = [shape for shape in new_slide.shapes if hasattr(shape, 'text_frame')]
        
        if len(textboxes) >= 1:
            # First textbox is heading
            textboxes[0].text_frame.text = heading_text
            self._apply_font_to_shape(textboxes[0])
        
        if len(textboxes) >= 2:
            # Second textbox is bullet points
            text_frame = textboxes[1].text_frame
            text_frame.clear()
            
            for bullet_point in bullet_points:
                p = text_frame.add_paragraph()
                p.text = bullet_point
            
            self._apply_font_to_shape(textboxes[1])
        
        # Add citations if provided
        self._add_citations_to_slide(new_slide, citations)
    
    def add_two_column_slide(self, heading_text: str, left_content: List[str], right_content: List[str], citations: List[str] = None):
        """Add two column slide (slide 3 from template) with heading and two columns."""
        # Copy slide 2 (two column slide) from template
        template_slide = self.template.slides[2]
        slide_layout = template_slide.slide_layout
        new_slide = self.presentation.slides.add_slide(slide_layout)
        
        # Find textboxes and update content
        textboxes = [shape for shape in new_slide.shapes if hasattr(shape, 'text_frame')]
        
        if len(textboxes) >= 1:
            # First textbox is heading
            textboxes[0].text_frame.text = heading_text
            self._apply_font_to_shape(textboxes[0])
        
        if len(textboxes) >= 2:
            # Second textbox is left column
            text_frame = textboxes[1].text_frame
            text_frame.clear()
            for item in left_content:
                p = text_frame.add_paragraph()
                p.text = item
            self._apply_font_to_shape(textboxes[1])
        
        if len(textboxes) >= 3:
            # Third textbox is right column
            text_frame = textboxes[2].text_frame
            text_frame.clear()
            for item in right_content:
                p = text_frame.add_paragraph()
                p.text = item
            self._apply_font_to_shape(textboxes[2])
        
        # Add citations if provided
        self._add_citations_to_slide(new_slide, citations)
    
    def add_image_slide(self, main_heading: str, sub_heading: str, citations: List[str] = None):
        """Add image slide (slide 4 from template) with main heading and sub heading."""
        # Copy slide 3 (image slide) from template
        template_slide = self.template.slides[3]
        slide_layout = template_slide.slide_layout
        new_slide = self.presentation.slides.add_slide(slide_layout)
        
        # Find textboxes and update content (image remains unchanged)
        textboxes = [shape for shape in new_slide.shapes if hasattr(shape, 'text_frame')]
        
        if len(textboxes) >= 1:
            # First textbox is main heading
            textboxes[0].text_frame.text = main_heading
            self._apply_font_to_shape(textboxes[0])
        
        if len(textboxes) >= 2:
            # Second textbox is sub heading
            textboxes[1].text_frame.text = sub_heading
            self._apply_font_to_shape(textboxes[1])
        
        # Add citations if provided
        self._add_citations_to_slide(new_slide, citations)
    
    def save(self):
        """Save the presentation."""
        self.presentation.save(self.output_filename)


if __name__ == "__main__":
    # Create presentation generator with custom font
    ppt = PPTGenerator("demo_presentation.pptx", font_name="Calibri")
    
    # Add title slide
    ppt.add_title_slide("My Demo Presentation")
    
    # Add bullet slide
    ppt.add_bullet_slide(
        "Key Features",
        [
            "Easy to use template system",
            "Preserves original design",
            "Only modifies text content",
            "Supports all four slide types"
        ]
    )
    
    # Add two column slide
    ppt.add_two_column_slide(
        "Comparison",
        [
            "Advantages:",
            "• Template-based",
            "• Consistent design",
            "• Fast creation"
        ],
        [
            "Features:",
            "• Title slides",
            "• Bullet points",
            "• Two columns",
            "• Image slides"
        ]
    )
    
    # Add image slide
    ppt.add_image_slide(
        "Summary",
        "This presentation was created using the template-based PPT generator"
    )

    ppt.add_bullet_slide(
        "Key Features",
        [
            "Easy to use template system",
            "Preserves original design",
            "Only modifies text content",
            "Supports all four slide types"
        ]
    )
    
    # Save presentation
    ppt.save()
    print("Demo presentation created: demo_presentation.pptx")
