"""
Advanced Template Management System for PPT Generation
Supports multiple templates, aspect ratios, and template customization
"""
import os
import logging
from typing import Dict, List, Optional, Tuple
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from django.conf import settings
import json
import hashlib
from dataclasses import dataclass, asdict
from enum import Enum
from .performance_cache import performance_cache, cached, cache_template_data

logger = logging.getLogger(__name__)


class AspectRatio(Enum):
    """Supported aspect ratios"""
    WIDESCREEN_16_9 = "16:9"
    STANDARD_4_3 = "4:3"
    WIDESCREEN_16_10 = "16:10"
    SQUARE_1_1 = "1:1"


class TemplateCategory(Enum):
    """Template categories"""
    BUSINESS = "business"
    ACADEMIC = "academic"
    CREATIVE = "creative"
    MEDICAL = "medical"
    TECHNOLOGY = "technology"
    EDUCATION = "education"


@dataclass
class TemplateMetadata:
    """Template metadata structure"""
    name: str
    category: TemplateCategory
    aspect_ratio: AspectRatio
    description: str
    author: str = "System"
    version: str = "1.0"
    slide_layouts: List[str] = None
    color_scheme: Dict[str, str] = None
    font_recommendations: List[str] = None
    created_at: str = None
    
    def __post_init__(self):
        if self.slide_layouts is None:
            self.slide_layouts = ["title", "bullet", "two-column", "content-image"]
        if self.color_scheme is None:
            self.color_scheme = {"primary": "#1f4e79", "secondary": "#70ad47", "accent": "#ffc000"}
        if self.font_recommendations is None:
            self.font_recommendations = ["Calibri", "Arial", "Segoe UI"]


class TemplateManager:
    """Advanced template management system"""
    
    def __init__(self):
        self.templates_dir = Path(settings.BASE_DIR) / "templates" / "presentations"
        self.metadata_dir = self.templates_dir / "metadata"
        self.cache_dir = self.templates_dir / "cache"
        
        # Create directories if they don't exist
        self.templates_dir.mkdir(parents=True, exist_ok=True)
        self.metadata_dir.mkdir(parents=True, exist_ok=True)
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        
        self._template_cache = {}
        self._metadata_cache = {}
        
        # Initialize default templates
        self._initialize_default_templates()
    
    def _initialize_default_templates(self):
        """Initialize default templates if they don't exist"""
        default_template_path = Path(settings.BASE_DIR) / "template.pptx"
        
        if default_template_path.exists():
            # Copy default template to templates directory
            default_dest = self.templates_dir / "default_16_9.pptx"
            if not default_dest.exists():
                import shutil
                shutil.copy2(default_template_path, default_dest)
                
                # Create metadata for default template
                metadata = TemplateMetadata(
                    name="Default Business Template",
                    category=TemplateCategory.BUSINESS,
                    aspect_ratio=AspectRatio.WIDESCREEN_16_9,
                    description="Standard business presentation template with clean design",
                    slide_layouts=["title", "bullet", "two-column", "content-image"]
                )
                self._save_template_metadata("default_16_9", metadata)
    
    def get_available_templates(self) -> Dict[str, TemplateMetadata]:
        """Get all available templates with their metadata"""
        templates = {}
        
        for template_file in self.templates_dir.glob("*.pptx"):
            template_id = template_file.stem
            metadata = self.get_template_metadata(template_id)
            if metadata:
                templates[template_id] = metadata
        
        return templates
    
    def get_templates_by_aspect_ratio(self, aspect_ratio: AspectRatio) -> Dict[str, TemplateMetadata]:
        """Get templates filtered by aspect ratio"""
        all_templates = self.get_available_templates()
        return {
            template_id: metadata 
            for template_id, metadata in all_templates.items()
            if metadata.aspect_ratio == aspect_ratio
        }
    
    def get_templates_by_category(self, category: TemplateCategory) -> Dict[str, TemplateMetadata]:
        """Get templates filtered by category"""
        all_templates = self.get_available_templates()
        return {
            template_id: metadata 
            for template_id, metadata in all_templates.items()
            if metadata.category == category
        }
    
    def get_template_metadata(self, template_id: str) -> Optional[TemplateMetadata]:
        """Get metadata for a specific template"""
        if template_id in self._metadata_cache:
            return self._metadata_cache[template_id]
        
        metadata_file = self.metadata_dir / f"{template_id}.json"
        if not metadata_file.exists():
            return None
        
        try:
            with open(metadata_file, 'r') as f:
                data = json.load(f)
            
            # Convert string enums back to enum objects
            data['category'] = TemplateCategory(data['category'])
            data['aspect_ratio'] = AspectRatio(data['aspect_ratio'])
            
            metadata = TemplateMetadata(**data)
            self._metadata_cache[template_id] = metadata
            return metadata
            
        except Exception as e:
            logger.error(f"Error loading template metadata for {template_id}: {e}")
            return None
    
    def _save_template_metadata(self, template_id: str, metadata: TemplateMetadata):
        """Save template metadata to file"""
        metadata_file = self.metadata_dir / f"{template_id}.json"
        
        try:
            # Convert to dict and handle enums
            data = asdict(metadata)
            data['category'] = metadata.category.value
            data['aspect_ratio'] = metadata.aspect_ratio.value
            
            with open(metadata_file, 'w') as f:
                json.dump(data, f, indent=2)
            
            # Update cache
            self._metadata_cache[template_id] = metadata
            
        except Exception as e:
            logger.error(f"Error saving template metadata for {template_id}: {e}")
    
    def load_template(self, template_id: str) -> Optional[Presentation]:
        """Load a template presentation with caching"""
        # Check performance cache first
        cached_template = performance_cache.get('template_data', template_id)
        if cached_template is not None:
            logger.debug(f"Using cached template: {template_id}")
            return cached_template
        
        # Check memory cache
        if template_id in self._template_cache:
            template = self._template_cache[template_id]
            # Promote to performance cache
            performance_cache.set('template_data', template_id, template, ttl=86400)
            return template
        
        template_file = self.templates_dir / f"{template_id}.pptx"
        if not template_file.exists():
            logger.error(f"Template file not found: {template_file}")
            return None
        
        try:
            presentation = Presentation(str(template_file))
            
            # Cache in both memory and performance cache
            self._template_cache[template_id] = presentation
            performance_cache.set('template_data', template_id, presentation, ttl=86400)
            
            logger.info(f"Loaded and cached template: {template_id}")
            return presentation
            
        except Exception as e:
            logger.error(f"Error loading template {template_id}: {e}")
            return None
    
    def get_template_dimensions(self, template_id: str) -> Optional[Tuple[int, int]]:
        """Get template dimensions (width, height)"""
        template = self.load_template(template_id)
        if not template:
            return None
        
        return (template.slide_width, template.slide_height)
    
    def validate_template(self, template_path: str) -> Tuple[bool, List[str]]:
        """Validate a template file"""
        errors = []
        
        try:
            presentation = Presentation(template_path)
            
            # Check minimum number of slides
            if len(presentation.slides) < 4:
                errors.append("Template must have at least 4 slides (title, bullet, two-column, content-image)")
            
            # Check slide layouts
            required_layouts = ["title", "bullet", "two-column", "content-image"]
            for i, layout_name in enumerate(required_layouts):
                if i >= len(presentation.slides):
                    errors.append(f"Missing {layout_name} slide at position {i}")
                    continue
                
                slide = presentation.slides[i]
                text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame')]
                
                if layout_name == "title" and len(text_shapes) < 1:
                    errors.append(f"Title slide must have at least 1 text shape")
                elif layout_name in ["bullet", "two-column"] and len(text_shapes) < 2:
                    errors.append(f"{layout_name} slide must have at least 2 text shapes")
                elif layout_name == "content-image" and len(text_shapes) < 2:
                    errors.append(f"Content-image slide must have at least 2 text shapes")
            
            return len(errors) == 0, errors
            
        except Exception as e:
            errors.append(f"Error validating template: {str(e)}")
            return False, errors
    
    def add_template(self, template_path: str, metadata: TemplateMetadata, template_id: str = None) -> Tuple[bool, str]:
        """Add a new template to the system"""
        try:
            # Validate template
            is_valid, errors = self.validate_template(template_path)
            if not is_valid:
                return False, f"Template validation failed: {'; '.join(errors)}"
            
            # Generate template ID if not provided
            if not template_id:
                template_id = self._generate_template_id(metadata.name)
            
            # Copy template file
            dest_path = self.templates_dir / f"{template_id}.pptx"
            import shutil
            shutil.copy2(template_path, dest_path)
            
            # Save metadata
            self._save_template_metadata(template_id, metadata)
            
            # Clear cache
            if template_id in self._template_cache:
                del self._template_cache[template_id]
            
            logger.info(f"Successfully added template: {template_id}")
            return True, template_id
            
        except Exception as e:
            logger.error(f"Error adding template: {e}")
            return False, str(e)
    
    def _generate_template_id(self, name: str) -> str:
        """Generate a unique template ID from name"""
        base_id = name.lower().replace(' ', '_').replace('-', '_')
        base_id = ''.join(c for c in base_id if c.isalnum() or c == '_')
        
        # Add hash to ensure uniqueness
        hash_suffix = hashlib.md5(name.encode()).hexdigest()[:8]
        return f"{base_id}_{hash_suffix}"
    
    def remove_template(self, template_id: str) -> bool:
        """Remove a template from the system"""
        try:
            template_file = self.templates_dir / f"{template_id}.pptx"
            metadata_file = self.metadata_dir / f"{template_id}.json"
            
            # Remove files
            if template_file.exists():
                template_file.unlink()
            if metadata_file.exists():
                metadata_file.unlink()
            
            # Clear cache
            if template_id in self._template_cache:
                del self._template_cache[template_id]
            if template_id in self._metadata_cache:
                del self._metadata_cache[template_id]
            
            logger.info(f"Successfully removed template: {template_id}")
            return True
            
        except Exception as e:
            logger.error(f"Error removing template {template_id}: {e}")
            return False
    
    def get_aspect_ratio_dimensions(self, aspect_ratio: AspectRatio) -> Tuple[int, int]:
        """Get standard dimensions for aspect ratio"""
        dimensions = {
            AspectRatio.WIDESCREEN_16_9: (Inches(13.333), Inches(7.5)),
            AspectRatio.STANDARD_4_3: (Inches(10), Inches(7.5)),
            AspectRatio.WIDESCREEN_16_10: (Inches(12), Inches(7.5)),
            AspectRatio.SQUARE_1_1: (Inches(7.5), Inches(7.5))
        }
        return dimensions.get(aspect_ratio, dimensions[AspectRatio.WIDESCREEN_16_9])
    
    def clear_cache(self):
        """Clear template cache"""
        self._template_cache.clear()
        self._metadata_cache.clear()
        logger.info("Template cache cleared")


# Global template manager instance
template_manager = TemplateManager()
